from pptx import Presentation
from pptx.util import Inches

def create_ppt_from_template(template_path, output_path, content):
    """
    Generate a PowerPoint presentation using a given template and content.

    :param template_path: Path to the PowerPoint template (.pptx).
    :param output_path: Path to save the generated PowerPoint file.
    :param content: List of dictionaries with 'title', 'subtitle' and 'body' for each slide.
    """
    # Load the PowerPoint template
    prs = Presentation(template_path)

    def get_placeholder_indices(slide):
        """Helper function to get placeholder indices for a specific slide"""
        indices = {}
        for shape in slide.placeholders:
            indices[shape.name] = {
                'idx': shape.placeholder_format.idx,
                'type': shape.name
            }
        return indices

    for slide_content in content:
        slide_layout = prs.slide_layouts[4]  # Title and Content with subtitle layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Get placeholder indices for this specific slide
        placeholders = get_placeholder_indices(slide)
        print(f"Placeholders for slide: {placeholders}")  # Debug info
        
        # Populate placeholders using indices
        for name, info in placeholders.items():
            shape = slide.placeholders[info['idx']]
            if "Title" in name:
                shape.text = slide_content.get("title", "")
            elif "Subtitle" in name:
                shape.text = slide_content.get("subtitle", "")
            elif "Content" in name or "Text" in name:
                shape.text = slide_content.get("body", "")

    # Save the new presentation
    prs.save(output_path)
    print(f"PPT saved at {output_path}")

# Example usage
template_path = "./Prismatic design.pptx"  # Path to an existing PowerPoint template
output_path = "output.pptx"  # Path to save the generated PowerPoint

content = [
    {
        "title": "Introduction",
        "subtitle": "Getting Started",
        "body": "Welcome to our presentation."
    },
    {
        "title": "Agenda",
        "subtitle": "Today's Topics",
        "body": "1. Overview\n2. Details\n3. Q&A"
    },
    {
        "title": "Conclusion",
        "subtitle": "Final Thoughts",
        "body": "Thank you for your time."
    }
]

create_ppt_from_template(template_path, output_path, content)

# Flow:
# Upload ppt template as .pptx
# Query each slide for placeholders and save the info as text
# Let user provide topic to generate content
# The content generated should be dependent on the placeholders given(pass both placeholders and topic to llm)
# Use Markers like #1# to specify which template slide to populate
#  return the ppt