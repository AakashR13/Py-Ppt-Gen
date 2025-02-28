from pptx import Presentation
import os
import requests
import streamlit as st
import json

from google import genai
from google.genai import types
from PIL import Image
from io import BytesIO

from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Get token from environment variable
LLMF_TOK = os.getenv('LLMF_TOK')

if not LLMF_TOK:
    raise ValueError("LLMF_TOK environment variable is not set")
def get_placeholder_indices(slide):
    """Helper function to get placeholder indices for a specific slide"""
    indices = {}
    for shape in slide.placeholders:
        indices[shape.name] = {
            'idx': shape.placeholder_format.idx,
            'type': shape.name
        }
    return indices

def describe_presentation(ppt_path):
    """
    Analyze all slides in a PowerPoint presentation and create descriptions
    Returns a list of dictionaries containing information about each slide
    """
    prs = Presentation(ppt_path)
    slides_info = []
    
    for slide_number, slide in enumerate(prs.slides, 1):
        # Get placeholder information
        placeholders = get_placeholder_indices(slide)
        
        # Get text content from all shapes
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
        
        # Create slide description
        slide_info = {
            'slide_number': slide_number,
            'layout_name': slide.slide_layout.name,
            'placeholders': placeholders,
            'text_content': text_content
        }
        slides_info.append(slide_info)
    
    return slides_info

# Replace the direct file handling with Streamlit interface
st.title("PowerPoint Analyzer")

# File uploader
uploaded_file = st.file_uploader("Choose a PowerPoint file", type=['pptx'])

if uploaded_file is not None:
    # Create description
    slides_description = describe_presentation(uploaded_file)
    
    # Build the description string
    description_text = ""
    for slide in slides_description:
        description_text += f"\nSlide {slide['slide_number']}:\n"
        description_text += f"Layout: {slide['layout_name']}\n"
        description_text += f"Placeholders: {slide['placeholders']}\n"
        description_text += f"Text content: {slide['text_content']}\n"
    
    # Display the description
    st.text_area("Presentation Description", description_text, height=300)
    
    # User input for content generation
    generation_prompt = st.text_input("Enter a prompt to generate presentation content:")
    
    if generation_prompt:
        # Create a prompt that asks for structured content based on available slides
        structured_prompt = f"""Based on the following presentation structure:
{description_text}

Generate presentation content for the prompt: {generation_prompt}
Return a JSON object with the following structure:
{{
    "slides": [
        {{
            "slide_number": integer,
            "layout_index": integer,  // Index of the slide layout (0-based)
            "content": {{
                // IMPORTANT: Use EXACTLY these placeholder names as shown below for each slide
            }}
        }}
    ]
}}

Important: 
1. Use ONLY the exact placeholder names provided for each slide.
2. Try to fill each text placeholder with 3-5 lines of content and match with a suitable layout.
3. Use a variety of layouts to populate with content.
4. Don't leave any placeholder key empty.
5. For any image/picture placeholders, provide a suitable caption for the image
"""
        
        client = genai.Client(api_key='AIzaSyCFyf_ljcTHvlN1jZdNyBTs_editn06Ggc')
        # Add available placeholders for each slide
        for slide in slides_description:
            structured_prompt += f"\nSlide {slide['slide_number']} placeholders: {', '.join([f'"{name}"' for name in slide['placeholders'].keys()])}"

        structured_prompt += "\n\n3. Do not make up new placeholder names"

        response = requests.post(
            "https://llmfoundry.straive.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {LLMF_TOK}"},
            json={
                "model": "gpt-4o-mini",
                "messages": [{"role": "user", "content": structured_prompt}],
                "response_format": {"type": "json_object"}
            }
        )

        # Display the structured response and generate PPT
        if response.status_code == 200:
            llm_response = response.json()
            response_content = llm_response['choices'][0]['message']['content']
            # Format the JSON response for better readability
            formatted_json = json.dumps(json.loads(response_content), indent=2)
            st.write("Generated Content:")
            st.code(formatted_json, language="json")
            
            try:
                # Parse the JSON response
                content_data = json.loads(response_content)
                
                # Create a new presentation using the template file
                prs = Presentation(uploaded_file)
                
                # Store the original slides' placeholder information before removing them
                original_slides_info = []
                for slide in prs.slides:
                    slide_info = {
                        'layout': slide.slide_layout,
                        'placeholders': get_placeholder_indices(slide)
                    }
                    original_slides_info.append(slide_info)
                    # st.write(f"Original slide info collected:", slide_info)
                
                # Remove all existing slides
                for _ in range(len(prs.slides._sldIdLst)):
                    rId = prs.slides._sldIdLst[0].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[0]
                
                # Now reconstruct slides using the original layout and placeholder information
                for slide_data in content_data.get('slides', []):
                    slide_number = slide_data.get('slide_number', 1) - 1  # Convert to 0-based index
                    
                    # Use the original slide info if available
                    if slide_number < len(original_slides_info):
                        original_info = original_slides_info[slide_number]
                        new_slide = prs.slides.add_slide(original_info['layout'])
                        placeholders = original_info['placeholders']
                    else:
                        # Fallback to layout_index if we don't have original info
                        layout_index = slide_data.get('layout_index', 0)
                        new_slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
                        placeholders = get_placeholder_indices(new_slide)
                    
                    st.write(f"Processing slide {slide_number + 1}")
                    st.write("Available placeholders:", placeholders)
                    
                    # Get content from formatted_json
                    content_dict = json.loads(formatted_json).get('slides', [{}])[slide_number].get('content', {})
                    
                    # Display available content for the slide with its respective placeholder data
                    # Use the content from the generated JSON instead of the actual placeholders
                    content_json = content_dict  # Use the generated content directly
                    st.json(content_json)  # Display content as JSON
                    
                    # Assign content to placeholders
                    for placeholder_name, content in content_json.items():
                        matched = False
                        for existing_name, existing_info in placeholders.items():
                            # Try to match the generated placeholder with existing ones
                            # Extract just the base name without numbers for matching
                            generated_base = ''.join([c for c in placeholder_name if not c.isdigit()]).strip()
                            existing_base = ''.join([c for c in existing_name if not c.isdigit()]).strip()
                            
                            if generated_base.lower() == existing_base.lower():  # Match base names
                                st.write(f"✅ Matched placeholder: '{placeholder_name}' with '{existing_name}'")
                                try:
                                    placeholder = new_slide.placeholders[existing_info['idx']]
                                    
                                    # Check if the placeholder is for an image
                                    if "picture" in existing_name.lower() and content != "No content available":
                                        st.write(f"Generating Picture: {content}")
                                        # Make a web request to generate the image
                                        response = client.models.generate_images(
                                            model='imagen-3.0-generate-002',
                                            prompt=content,
                                            config=types.GenerateImagesConfig(
                                                number_of_images=1,
                                            )
                                        )
                                        # Get the first (and only) generated image
                                        if response.generated_images:
                                            image_bytes = response.generated_images[0].image.image_bytes
                                            # Create a temporary file to save the image
                                            temp_img_path = f"temp_image_{slide_number}.png"
                                            
                                            # Save bytes directly to file
                                            with open(temp_img_path, "wb") as img_file:
                                                img_file.write(image_bytes)
                                            
                                            try:
                                                # Insert the image into the placeholder
                                                placeholder.insert_picture(temp_img_path)
                                                # Clean up the temporary file
                                                os.remove(temp_img_path)
                                            except Exception as e:
                                                st.error(f"Error inserting image: {e}")
                                        else:
                                            st.warning("No image was generated")
                                    else:
                                        # Handle text content that might be a dictionary
                                        if isinstance(content, dict) and 'text_content' in content:
                                            # Join multiple text items with newlines if it's a list
                                            if isinstance(content['text_content'], list):
                                                text_content = '\n'.join(str(item) for item in content['text_content'])
                                            else:
                                                text_content = str(content['text_content'])
                                            placeholder.text = text_content
                                        else:
                                            # Handle direct string content
                                            placeholder.text = str(content)
                                    matched = True
                                except KeyError:
                                    st.warning(f"Placeholder index {existing_info['idx']} not found on slide {slide_number + 1}. Skipping.")
                                except Exception as e:
                                    st.warning(f"Error setting content for placeholder '{existing_name}' on slide {slide_number + 1}: {str(e)}")
                                break
                        if not matched:
                            st.warning(f"Placeholder '{placeholder_name}' not found in slide {slide_number + 1}. Skipping assignment.")
                
                # Save and provide download
                output = "generated_presentation.pptx"
                prs.save(output)
                
                with open(output, "rb") as file:
                    st.download_button(
                        label="Download Generated Presentation",
                        data=file,
                        file_name="generated_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
            except json.JSONDecodeError as e:
                st.error(f"Error parsing LLM response: {e}")
            except Exception as e:
                st.error(f"Error generating presentation: {e}")
        else:
            st.error("Error getting response from LLM") 